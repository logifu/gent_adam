import docx
import openpyxl
import os
import sys

docs_dir = r"d:\HSIIB\docs"
out_dir = r"d:\HSIIB\docs\extracted"
os.makedirs(out_dir, exist_ok=True)

# Extract all docx files
docx_files = [f for f in os.listdir(docs_dir) if f.endswith('.docx')]
for fname in sorted(docx_files):
    fpath = os.path.join(docs_dir, fname)
    outpath = os.path.join(out_dir, fname.replace('.docx', '.txt'))
    try:
        doc = docx.Document(fpath)
        lines = []
        for para in doc.paragraphs:
            lines.append(para.text)
        # Also extract tables
        for i, table in enumerate(doc.tables):
            lines.append(f"\n--- TABLE {i+1} ---")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append(" | ".join(cells))
        with open(outpath, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        print(f"OK: {fname} -> {len(lines)} lines")
    except Exception as e:
        print(f"ERR: {fname}: {e}")

# Extract xlsx files
xlsx_files = [f for f in os.listdir(docs_dir) if f.endswith('.xlsx')]
for fname in sorted(xlsx_files):
    fpath = os.path.join(docs_dir, fname)
    outpath = os.path.join(out_dir, fname.replace('.xlsx', '.txt'))
    try:
        wb = openpyxl.load_workbook(fpath, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"\n=== SHEET: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) if c is not None else '' for c in row]
                lines.append(" | ".join(vals))
        with open(outpath, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        print(f"OK: {fname} -> {len(lines)} lines")
    except Exception as e:
        print(f"ERR: {fname}: {e}")

# Extract pptx files
try:
    from pptx import Presentation
    pptx_files = [f for f in os.listdir(docs_dir) if f.endswith('.pptx')]
    for fname in sorted(pptx_files):
        fpath = os.path.join(docs_dir, fname)
        outpath = os.path.join(out_dir, fname.replace('.pptx', '.txt'))
        try:
            prs = Presentation(fpath)
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"\n=== SLIDE {i+1} ===")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            lines.append(para.text)
            with open(outpath, 'w', encoding='utf-8') as f:
                f.write('\n'.join(lines))
            print(f"OK: {fname} -> {len(lines)} lines")
        except Exception as e:
            print(f"ERR: {fname}: {e}")
except ImportError:
    print("pptx module not available, skipping pptx files")

print("\nDone! Files written to:", out_dir)
